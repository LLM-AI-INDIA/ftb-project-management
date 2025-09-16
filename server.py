#!/usr/bin/env python3
"""
Google Drive MCP Server
Reads files from Google Drive and provides content for analysis
"""

import asyncio
import os
import io
import logging
from pptx import Presentation          # pip install python-pptx
import docx                            # pip install python-docx
import xml.etree.ElementTree as ET
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import pandas as pd
import PyPDF2
from typing import Any, Dict, List, Optional
import json

# MCP imports
from mcp.server.models import InitializationOptions
from mcp.server import NotificationOptions, Server
from mcp.types import Resource, Tool, TextContent, ImageContent, EmbeddedResource
import mcp.types as types

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("gdrive-mcp-server")

# Google Drive API configuration - Updated scopes
SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/drive.metadata.readonly'
]

class GoogleDriveMCPServer:
    def __init__(self):
        self.service = None
        self.file_cache = {}
        
    def get_credentials(self):
        """Get Google Drive API credentials with better error handling"""
        # Use the directory where this script is located
        script_dir = os.path.dirname(os.path.abspath(__file__))
        token_path = os.path.join(script_dir, 'token.json')
        credentials_path = os.path.join(script_dir, 'credentials.json')
        
        creds = None
        
        # Try to load existing credentials
        if os.path.exists(token_path):
            try:
                creds = Credentials.from_authorized_user_file(token_path, SCOPES)
                logger.info("Loaded existing credentials from token.json")
            except Exception as e:
                logger.warning(f"Failed to load existing credentials: {e}")
                # Delete the invalid token file
                os.remove(token_path)
                creds = None
        
        # Check if credentials are valid
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                try:
                    logger.info("Refreshing expired credentials...")
                    creds.refresh(Request())
                    logger.info("Successfully refreshed credentials")
                except Exception as e:
                    logger.error(f"Failed to refresh credentials: {e}")
                    # Delete the invalid token file and re-authenticate
                    if os.path.exists(token_path):
                        os.remove(token_path)
                    creds = None
            
            # If we still don't have valid credentials, run the OAuth flow
            if not creds or not creds.valid:
                if not os.path.exists(credentials_path):
                    raise FileNotFoundError(
                        f"credentials.json not found at {credentials_path}. "
                        f"Please download it from Google Cloud Console and place it in {script_dir}"
                    )
                
                logger.info("Starting OAuth flow for new credentials...")
                flow = InstalledAppFlow.from_client_secrets_file(
                    credentials_path, SCOPES)
                creds = flow.run_local_server(port=0)
                logger.info("Successfully obtained new credentials")
            
            # Save the credentials for next run
            try:
                with open(token_path, 'w') as token:
                    token.write(creds.to_json())
                logger.info("Saved credentials to token.json")
            except Exception as e:
                logger.warning(f"Failed to save credentials: {e}")
        
        return creds

    def initialize_drive_service(self):
        """Initialize Google Drive service with better error handling"""
        try:
            logger.info("Initializing Google Drive service...")
            creds = self.get_credentials()
            self.service = build('drive', 'v3', credentials=creds)
            
            # Test the service by making a simple API call
            try:
                test_result = self.service.files().list(pageSize=1).execute()
                logger.info("Google Drive service initialized and tested successfully")
            except Exception as e:
                logger.error(f"Failed to test Google Drive service: {e}")
                raise
                
        except Exception as e:
            logger.error(f"Failed to initialize Google Drive service: {e}")
            raise

    def list_files(self, query: str = None, max_results: int = 100) -> List[Dict]:
        """List files from Google Drive"""
        if not self.service:
            self.initialize_drive_service()
        
        try:
            # Build query for supported file types
            supported_types = [
                "mimeType='text/csv'",
                "mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'",
                "mimeType='application/vnd.ms-project'",
                "mimeType='application/pdf'",
                "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'",
                "mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'",
                "mimeType='text/plain'",
                "mimeType='text/xml'",
                "mimeType='application/json'"
            ]
            
            base_query = f"({' or '.join(supported_types)})"
            if query:
                full_query = f"{base_query} and name contains '{query}'"
            else:
                full_query = base_query
            
            logger.info(f"Searching files with query: {full_query}")
            
            results = self.service.files().list(
                q=full_query,
                pageSize=max_results,
                fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)"
            ).execute()
            
            files = results.get('files', [])
            logger.info(f"Found {len(files)} files")
            return files
        
        except Exception as e:
            logger.error(f"Error listing files: {e}")
            return []

    def download_file(self, file_id: str) -> tuple:
        """Download file content from Google Drive"""
        if not self.service:
            self.initialize_drive_service()
        
        try:
            # Get file metadata
            file = self.service.files().get(fileId=file_id).execute()
            file_name = file['name']
            mime_type = file['mimeType']
            
            logger.info(f"Downloading file: {file_name} ({mime_type})")
            
            # Download file content
            request = self.service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            logger.info(f"Successfully downloaded file: {file_name}")
            return fh, file_name, mime_type
        
        except Exception as e:
            logger.error(f"Error downloading file {file_id}: {e}")
            raise

    def process_csv(self, fh: io.BytesIO) -> Dict:
        """Process CSV file - COMPLETE data extraction"""
        fh.seek(0)
        try:
            df = pd.read_csv(fh)
            return {
                "type": "csv",
                "shape": df.shape,
                "columns": df.columns.tolist(),
                "data": df.to_dict('records'),  # REMOVE the [:100] limit - get ALL data
                "summary": df.describe().to_dict() if len(df.select_dtypes(include=['number']).columns) > 0 else None,
                "full_data_extracted": True,
                "total_records": len(df)
            }
        except Exception as e:
            return {"type": "csv", "error": f"Failed to process CSV: {str(e)}"}

    def process_excel(self, fh: io.BytesIO) -> Dict:
        """Process Excel file - COMPLETE data extraction"""
        fh.seek(0)
        try:
            df = pd.read_excel(fh, engine='openpyxl')
            return {
                "type": "excel",
                "shape": df.shape,
                "columns": df.columns.tolist(),
                "data": df.to_dict('records'),  # REMOVE the [:100] limit - get ALL data
                "summary": df.describe().to_dict() if len(df.select_dtypes(include=['number']).columns) > 0 else None,
                "full_data_extracted": True,
                "total_records": len(df)
            }
        except Exception as e:
            return {"type": "excel", "error": f"Failed to process Excel: {str(e)}"}

    def process_pdf(self, fh: io.BytesIO) -> Dict:
        """Process PDF file - COMPLETE text extraction"""
        fh.seek(0)
        try:
            pdf_reader = PyPDF2.PdfReader(fh)
            text_content = ""

            # Extract from ALL pages, not just first 10
            for page_num, page in enumerate(pdf_reader.pages):
                text_content += f"--- Page {page_num + 1} ---\n"
                text_content += page.extract_text()
                text_content += "\n\n"

            return {
                "type": "pdf",
                "num_pages": len(pdf_reader.pages),
                "content": text_content,  # REMOVE the [:10000] limit - get ALL content
                "metadata": pdf_reader.metadata if hasattr(pdf_reader, 'metadata') else None,
                "full_content_extracted": True
            }
        except Exception as e:
            return {"type": "pdf", "error": f"Failed to process PDF: {str(e)}"}

    def process_docx(self, fh: io.BytesIO) -> Dict:
        """Process DOCX file - COMPLETE content extraction"""
        fh.seek(0)
        try:
            doc = docx.Document(fh)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            body = "\n".join(full_text)
            return {
                "type": "docx", 
                "content": body,  # REMOVE the [:10_000] limit - get ALL content
                "full_content_extracted": True,
                "total_paragraphs": len(full_text)
            }
        except Exception as e:
            return {"type": "docx", "error": str(e)}

    def process_pptx(self, fh: io.BytesIO) -> Dict:
        """Process PPTX file - COMPLETE slide extraction"""
        fh.seek(0)
        try:
            prs = Presentation(fh)
            slides = []
            # Process ALL slides, not just first 20
            for idx, s in enumerate(prs.slides):
                text = []
                for shape in s.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text.strip())
                slides.append({"slide": idx+1, "content": " | ".join(text)})
            return {
                "type": "pptx", 
                "num_slides": len(prs.slides), 
                "slides": slides,
                "full_content_extracted": True
            }
        except Exception as e:
            return {"type": "pptx", "error": str(e)}

    def process_txt(self, fh: io.BytesIO) -> Dict:
        fh.seek(0)
        try:
            text = fh.read().decode("utf-8")
            return {"type": "txt", "content": text[:10_000]}
        except Exception as e:
            return {"type": "txt", "error": str(e)}

    def process_xml(self, fh: io.BytesIO) -> Dict:
        fh.seek(0)
        try:
            root = ET.parse(fh).getroot()
            nodes = list(root.iter())
            preview = [f"<{el.tag}>{el.text.strip() if el.text else ''}</{el.tag}>"
                       for el in nodes[:500]]
            return {"type": "xml", "root_tag": root.tag, "nodes_count": len(nodes),
                    "preview": preview}
        except Exception as e:
            return {"type": "xml", "error": str(e)}

    def process_json(self, fh: io.BytesIO) -> Dict:
        """Process JSON file - COMPLETE content extraction"""
        fh.seek(0)
        try:
            data = json.load(fh)
            # Return complete JSON, not truncated
            out = json.dumps(data, indent=2)
            return {
                "type": "json", 
                "content": out,  # REMOVE the [:50_000] limit
                "full_content_extracted": True,
                "data_size": len(out)
            }
        except Exception as e:
            return {"type": "json", "error": str(e)}
    


    # STEP 2: Smart Chunking Strategy for Large Files
    def create_smart_chunks(content, file_name, file_type, max_chunk_size=8000):
        """Break large content into meaningful chunks"""

        if len(content) <= max_chunk_size:
            return [{"chunk": 1, "content": content, "is_complete": True}]

        chunks = []

        if file_type in ["csv", "excel"]:
            # For data files, chunk by records
            lines = content.split('\n')
            current_chunk = ""
            chunk_num = 1
            records_in_chunk = 0

            for line in lines:
                if len(current_chunk + line) > max_chunk_size and current_chunk:
                    chunks.append({
                        "chunk": chunk_num,
                        "content": current_chunk,
                        "records_count": records_in_chunk,
                        "is_complete": False
                    })
                    current_chunk = line + '\n'
                    chunk_num += 1
                    records_in_chunk = 1
                else:
                    current_chunk += line + '\n'
                    if line.strip():
                        records_in_chunk += 1

            if current_chunk:
                chunks.append({
                    "chunk": chunk_num,
                    "content": current_chunk,
                    "records_count": records_in_chunk,
                    "is_complete": chunk_num == 1
                })

        elif file_type == "pdf":
            # For PDFs, chunk by pages
            pages = content.split("--- Page ")
            current_chunk = ""
            chunk_num = 1

            for page in pages:
                if page.strip():
                    page_content = "--- Page " + page
                    if len(current_chunk + page_content) > max_chunk_size and current_chunk:
                        chunks.append({
                            "chunk": chunk_num,
                            "content": current_chunk,
                            "is_complete": False
                        })
                        current_chunk = page_content
                        chunk_num += 1
                    else:
                        current_chunk += page_content

            if current_chunk:
                chunks.append({
                    "chunk": chunk_num,
                    "content": current_chunk,
                    "is_complete": False
                })

        else:
            # For other files, chunk by sentences/paragraphs
            sentences = content.split('. ')
            current_chunk = ""
            chunk_num = 1

            for sentence in sentences:
                if len(current_chunk + sentence) > max_chunk_size and current_chunk:
                    chunks.append({
                        "chunk": chunk_num,
                        "content": current_chunk,
                        "is_complete": False
                    })
                    current_chunk = sentence + '. '
                    chunk_num += 1
                else:
                    current_chunk += sentence + '. '

            if current_chunk:
                chunks.append({
                    "chunk": chunk_num,
                    "content": current_chunk,
                    "is_complete": False
                })

        return chunks

    # STEP 3: Enhanced Flask app with complete processing
    def process_all_files_completely(http_server_url, max_files=100):
        """Process ALL files with complete content extraction"""

        print("Starting complete file processing...")

        # Get ALL files
        response = requests.post(f"{http_server_url}/call_tool", json={
            "name": "search_gdrive_files",
            "arguments": {"query": "", "max_results": max_files}
        })

        if response.status_code != 200:
            return {"success": False, "error": "Failed to connect to Google Drive"}

        result = response.json()
        all_files = result.get("data", []) if result.get("success") else []

        print(f"Processing ALL {len(all_files)} files completely...")

        complete_file_data = []
        processing_stats = {
            "total_files": len(all_files),
            "successfully_processed": 0,
            "failed_files": [],
            "large_files_chunked": 0,
            "total_chunks": 0
        }

        for i, file_info in enumerate(all_files):
            file_id = file_info["id"]
            file_name = file_info["name"]

            print(f"Processing {i+1}/{len(all_files)}: {file_name}")

            try:
                # Read complete file content
                response = requests.post(f"{http_server_url}/call_tool", json={
                    "name": "read_gdrive_file",
                    "arguments": {"file_id": file_id}
                })

                if response.status_code == 200:
                    result = response.json()
                    if result.get("success"):
                        file_data = result.get("data", {})
                        content_info = file_data.get("content", {})

                        # Extract COMPLETE content without truncation
                        complete_content = extract_complete_content(content_info, file_name)

                        if complete_content:
                            # Create smart chunks for large files
                            chunks = create_smart_chunks(
                                complete_content, 
                                file_name, 
                                content_info.get("type", "unknown")
                            )

                            complete_file_data.append({
                                "file_name": file_name,
                                "file_id": file_id,
                                "file_type": content_info.get("type", "unknown"),
                                "mime_type": file_info.get("mimeType", "unknown"),
                                "chunks": chunks,
                                "total_chunks": len(chunks),
                                "is_large_file": len(chunks) > 1,
                                "complete_content_length": len(complete_content)
                            })

                            processing_stats["successfully_processed"] += 1
                            if len(chunks) > 1:
                                processing_stats["large_files_chunked"] += 1
                            processing_stats["total_chunks"] += len(chunks)

                            print(f"✓ {file_name}: {len(chunks)} chunks, {len(complete_content)} characters")
                        else:
                            processing_stats["failed_files"].append(f"{file_name} (no content)")
                            print(f"✗ {file_name}: No content extracted")
                    else:
                        processing_stats["failed_files"].append(f"{file_name} (read failed)")
                        print(f"✗ {file_name}: Read failed")
                else:
                    processing_stats["failed_files"].append(f"{file_name} (HTTP error)")
                    print(f"✗ {file_name}: HTTP error")

            except Exception as e:
                processing_stats["failed_files"].append(f"{file_name} (exception)")
                print(f"✗ {file_name}: Exception - {str(e)}")

        print(f"\nCOMPLETE PROCESSING SUMMARY:")
        print(f"Total files: {processing_stats['total_files']}")
        print(f"Successfully processed: {processing_stats['successfully_processed']}")
        print(f"Failed: {len(processing_stats['failed_files'])}")
        print(f"Large files chunked: {processing_stats['large_files_chunked']}")
        print(f"Total chunks created: {processing_stats['total_chunks']}")

        return {
            "success": True,
            "file_data": complete_file_data,
            "stats": processing_stats
        }

    def extract_complete_content(content_info, file_name):
        """Extract complete content without any truncation"""

        file_type = content_info.get("type", "unknown")

        if file_type == "pdf":
            return content_info.get("content", "")
        elif file_type in ["csv", "excel"]:
            data = content_info.get("data", [])
            columns = content_info.get("columns", [])
            shape = content_info.get("shape", [0, 0])

            content = f"COMPLETE DATASET: {shape[0]} rows, {shape[1]} columns\n"
            content += f"Columns: {', '.join(columns)}\n\n"

            # Include ALL records, not just a sample
            content += "ALL RECORDS:\n"
            for i, record in enumerate(data, 1):
                record_text = []
                for k, v in record.items():
                    if v is not None:
                        record_text.append(f"{k}: {v}")
                content += f"Record {i}: {', '.join(record_text)}\n"

            return content
        elif file_type in ["docx", "txt"]:
            return content_info.get("content", "")
        elif file_type == "pptx":
            slides = content_info.get("slides", [])
            content = f"COMPLETE PRESENTATION: {len(slides)} slides\n\n"
            for slide in slides:
                content += f"Slide {slide.get('slide', '')}: {slide.get('content', '')}\n"
            return content
        elif file_type == "json":
            return content_info.get("content", "")
        elif file_type == "xml":
            preview = content_info.get("preview", [])
            return "\n".join(preview)

        return ""

        
    def process_file_content(self, fh: io.BytesIO, mime_type: str) -> Dict:
        """Process file content based on MIME type"""
        if mime_type == 'text/csv':
            return self.process_csv(fh)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return self.process_excel(fh)
        elif mime_type == 'application/pdf':
            return self.process_pdf(fh)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return self.process_pptx(fh)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self.process_docx(fh)
        elif mime_type == 'text/plain':
            return self.process_txt(fh)
        elif mime_type == 'text/xml':
            return self.process_xml(fh)
        elif mime_type == 'application/json':
            return self.process_json(fh)
        else:
            return {"type": "unknown", "error": f"Unsupported MIME type: {mime_type}"}

    async def read_file(self, file_id: str) -> Dict:
        """Read and process file from Google Drive"""
        try:
            # Check cache first
            if file_id in self.file_cache:
                logger.info(f"Returning cached result for file {file_id}")
                return self.file_cache[file_id]
            
            # Download and process file
            fh, file_name, mime_type = self.download_file(file_id)
            content = self.process_file_content(fh, mime_type)
            
            result = {
                "file_id": file_id,
                "file_name": file_name,
                "mime_type": mime_type,
                "content": content
            }
            
            # Cache the result
            self.file_cache[file_id] = result
            logger.info(f"Processed and cached file: {file_name}")
            
            return result
        
        except Exception as e:
            logger.error(f"Failed to read file {file_id}: {e}")
            return {"error": f"Failed to read file: {str(e)}"}


# Initialize the MCP server
app = Server("gdrive-mcp-server")
gdrive_server = GoogleDriveMCPServer()





def list_folders(self, parent_folder_id='root', max_results=100):
    """List folders in Google Drive"""
    if not self.service:
        self.initialize_drive_service()
    
    try:
        # Query to find folders
        folder_query = "mimeType='application/vnd.google-apps.folder'"
        if parent_folder_id != 'root':
            folder_query += f" and '{parent_folder_id}' in parents"
        
        results = self.service.files().list(
            q=folder_query,
            pageSize=max_results,
            fields="nextPageToken, files(id, name, parents, modifiedTime)"
        ).execute()
        
        folders = results.get('files', [])
        logger.info(f"Found {len(folders)} folders")
        return folders
    
    except Exception as e:
        logger.error(f"Error listing folders: {e}")
        return []

def list_files_in_folder(self, folder_id, max_results=100):
    """List files within a specific folder"""
    if not self.service:
        self.initialize_drive_service()
    
    try:
        # Build query for supported file types within the folder
        supported_types = [
            "mimeType='text/csv'",
            "mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'",
            "mimeType='application/vnd.ms-project'",
            "mimeType='application/pdf'",
            "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'",
            "mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'",
            "mimeType='text/plain'",
            "mimeType='text/xml'",
            "mimeType='application/json'"
        ]
        
        file_query = f"({' or '.join(supported_types)}) and '{folder_id}' in parents"
        
        results = self.service.files().list(
            q=file_query,
            pageSize=max_results,
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime, parents)"
        ).execute()
        
        files = results.get('files', [])
        logger.info(f"Found {len(files)} files in folder {folder_id}")
        return files
    
    except Exception as e:
        logger.error(f"Error listing files in folder {folder_id}: {e}")
        return []

def get_folder_structure(self, max_depth=2):
    """Get the complete folder structure"""
    if not self.service:
        self.initialize_drive_service()
    
    try:
        folder_structure = {
            'root': {
                'name': 'My Drive',
                'id': 'root',
                'subfolders': [],
                'files': []
            }
        }
        
        def explore_folder(folder_id, folder_name, current_depth=0):
            if current_depth >= max_depth:
                return []
            
            # Get subfolders
            subfolders = self.list_folders(folder_id)
            subfolder_data = []
            
            for subfolder in subfolders:
                subfolder_info = {
                    'name': subfolder['name'],
                    'id': subfolder['id'],
                    'subfolders': explore_folder(subfolder['id'], subfolder['name'], current_depth + 1),
                    'files': self.list_files_in_folder(subfolder['id'])
                }
                subfolder_data.append(subfolder_info)
            
            return subfolder_data
        
        # Get root level files and folders
        folder_structure['root']['files'] = self.list_files()  # Your existing method
        folder_structure['root']['subfolders'] = explore_folder('root', 'My Drive')
        
        return folder_structure
    
    except Exception as e:
        logger.error(f"Error getting folder structure: {e}")
        return {}

def find_folder_by_name(self, folder_name):
    """Find folders by name"""
    if not self.service:
        self.initialize_drive_service()
    
    try:
        query = f"mimeType='application/vnd.google-apps.folder' and name contains '{folder_name}'"
        
        results = self.service.files().list(
            q=query,
            pageSize=20,
            fields="nextPageToken, files(id, name, parents, modifiedTime)"
        ).execute()
        
        folders = results.get('files', [])
        logger.info(f"Found {len(folders)} folders matching '{folder_name}'")
        return folders
    
    except Exception as e:
        logger.error(f"Error finding folder '{folder_name}': {e}")
        return []



    

@app.list_resources()
async def handle_list_resources() -> list[Resource]:
    """List available Google Drive files as resources"""
    files = gdrive_server.list_files()
    resources = []
    
    for file in files:
        resources.append(Resource(
            uri=f"gdrive://{file['id']}",
            name=file['name'],
            description=f"Google Drive file: {file['name']} ({file['mimeType']})",
            mimeType=file.get('mimeType', 'application/octet-stream')
        ))
    
    return resources


@app.read_resource()
async def handle_read_resource(uri: str) -> str:
    """Read a specific resource (file) from Google Drive"""
    if not uri.startswith("gdrive://"):
        raise ValueError(f"Invalid URI: {uri}")
    
    file_id = uri.replace("gdrive://", "")
    result = await gdrive_server.read_file(file_id)
    
    return json.dumps(result, indent=2, default=str)


@app.list_tools()
async def handle_list_tools() -> list[Tool]:
    """List available tools"""
    return [
        Tool(
            name="search_gdrive_files",
            description="Search for files in Google Drive by name",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Search query for file names"
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Maximum number of results to return",
                        "default": 10
                    }
                },
                "required": ["query"]
            }
        ),
        Tool(
            name="read_gdrive_file",
            description="Read and analyze a specific Google Drive file by ID",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_id": {
                        "type": "string",
                        "description": "Google Drive file ID"
                    }
                },
                "required": ["file_id"]
            }
        ),
        Tool(
            name="list_folders",
            description="List folders in Google Drive",
            inputSchema={
                "type": "object",
                "properties": {
                    "parent_folder_id": {"type": "string", "description": "Parent folder ID (use 'root' for root directory)", "default": "root"},
                    "max_results": {"type": "integer", "description": "Maximum number of results", "default": 100}
                }
            }
        ),
        Tool(
            name="find_folder",
            description="Find folders by name",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_name": {"type": "string", "description": "Name of folder to find"}
                },
                "required": ["folder_name"]
            }
        ),
        Tool(
            name="list_files_in_folder",
            description="List files within a specific folder",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_id": {"type": "string", "description": "Folder ID to list files from"}
                },
                "required": ["folder_id"]
            }
        ),
        Tool(
            name="get_folder_structure",
            description="Get complete folder structure",
            inputSchema={
                "type": "object",
                "properties": {
                    "max_depth": {"type": "integer", "description": "Maximum depth to explore", "default": 2}
                }
            }
        ),        
        Tool(
            name="analyze_file_data",
            description="Analyze data from a previously read file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_id": {
                        "type": "string",
                        "description": "Google Drive file ID to analyze"
                    },
                    "analysis_type": {
                        "type": "string",
                        "enum": ["summary", "statistics", "structure", "content"],
                        "description": "Type of analysis to perform"
                    }
                },
                "required": ["file_id", "analysis_type"]
            }
        )
    ]


@app.call_tool()
async def handle_call_tool(name: str, arguments: dict) -> list[types.TextContent]:
    """Handle tool calls"""
    try:
        if name == "list_folders":
            parent_folder_id = arguments.get("parent_folder_id", "root")
            max_results = arguments.get("max_results", 100)
            
            folders = gdrive_server.list_folders(parent_folder_id, max_results)
            
            result = f"Folders in {'root directory' if parent_folder_id == 'root' else parent_folder_id}:\n"
            for folder in folders:
                result += f"- {folder['name']} (ID: {folder['id']})\n"
            
            return [types.TextContent(type="text", text=result)]
        
        elif name == "find_folder":
            folder_name = arguments.get("folder_name")
            if not folder_name:
                return [types.TextContent(type="text", text="Error: folder_name is required")]
            
            folders = gdrive_server.find_folder_by_name(folder_name)
            
            result = f"Folders matching '{folder_name}':\n"
            for folder in folders:
                result += f"- {folder['name']} (ID: {folder['id']})\n"
            
            return [types.TextContent(type="text", text=result)]
        
        elif name == "list_files_in_folder":
            folder_id = arguments.get("folder_id")
            if not folder_id:
                return [types.TextContent(type="text", text="Error: folder_id is required")]
            
            files = gdrive_server.list_files_in_folder(folder_id)
            
            result = f"Files in folder {folder_id}:\n"
            for file in files:
                result += f"- {file['name']} (ID: {file['id']}, Type: {file['mimeType']})\n"
            
            return [types.TextContent(type="text", text=result)]
        
        elif name == "get_folder_structure":
            max_depth = arguments.get("max_depth", 2)
            
            structure = gdrive_server.get_folder_structure(max_depth)
            
            return [types.TextContent(type="text", text=json.dumps(structure, indent=2, default=str))]
        
        elif name == "search_gdrive_files":
            query = arguments.get("query", "")
            max_results = arguments.get("max_results", 10)
            
            files = gdrive_server.list_files(query=query, max_results=max_results)
            
            result = "Found files:\n"
            for file in files:
                result += f"- {file['name']} (ID: {file['id']}, Type: {file['mimeType']})\n"
            
            return [types.TextContent(type="text", text=result)]
        
        elif name == "read_gdrive_file":
            file_id = arguments.get("file_id")
            if not file_id:
                return [types.TextContent(type="text", text="Error: file_id is required")]
            
            result = await gdrive_server.read_file(file_id)
            return [types.TextContent(type="text", text=json.dumps(result, indent=2, default=str))]
        
        elif name == "analyze_file_data":
            file_id = arguments.get("file_id")
            analysis_type = arguments.get("analysis_type")
            
            if not file_id or not analysis_type:
                return [types.TextContent(type="text", text="Error: file_id and analysis_type are required")]
            
            # Read the file first
            file_data = await gdrive_server.read_file(file_id)
            
            if "error" in file_data:
                return [types.TextContent(type="text", text=f"Error reading file: {file_data['error']}")]
            
            content = file_data.get("content", {})
            file_type = content.get("type", "unknown")
            
            analysis_result = ""
            
            if analysis_type == "summary":
                analysis_result = f"File Summary for {file_data['file_name']}:\n"
                analysis_result += f"Type: {file_type}\n"
                if file_type in ["csv", "excel"]:
                    analysis_result += f"Shape: {content.get('shape', 'Unknown')}\n"
                    analysis_result += f"Columns: {', '.join(content.get('columns', []))}\n"
                elif file_type == "pdf":
                    analysis_result += f"Pages: {content.get('num_pages', 'Unknown')}\n"

            
            elif analysis_type == "statistics":
                if file_type in ["csv", "excel"] and content.get("summary"):
                    analysis_result = f"Statistical Summary:\n{json.dumps(content['summary'], indent=2)}"
                else:
                    analysis_result = "Statistics not available for this file type"
            
            elif analysis_type == "structure":
                analysis_result = f"File Structure:\n{json.dumps(content, indent=2, default=str)}"
            
            elif analysis_type == "content":
                if file_type == "pdf":
                    analysis_result = f"PDF Content Preview:\n{content.get('content', 'No content available')}"
                elif file_type in ["csv", "excel"]:
                    analysis_result = f"Data Preview (first few rows):\n{json.dumps(content.get('data', [])[:5], indent=2)}"
                else:
                    analysis_result = "Content preview not available for this file type"
            
            return [types.TextContent(type="text", text=analysis_result)]
        
        else:
            return [types.TextContent(type="text", text=f"Unknown tool: {name}")]
    
    except Exception as e:
        logger.error(f"Error in tool call {name}: {e}")
        return [types.TextContent(type="text", text=f"Error: {str(e)}")]


async def main():
    # Import here to avoid issues if mcp package is not available
    from mcp.server.stdio import stdio_server
    
    async with stdio_server() as (read_stream, write_stream):
        logger.info("MCP server is starting...")
        await app.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="gdrive-mcp-server",
                server_version="0.1.0",
                capabilities=app.get_capabilities(
                    notification_options=NotificationOptions(),
                    experimental_capabilities={},
                ),
            ),
        )
        logger.info("MCP server has started successfully.")

if __name__ == "__main__":
    # This should only be used for the HTTP wrapper, not directly
    logger.error("This script should be imported by http_server.py, not run directly")
